#!/usr/bin/env python3
"""Generate product introduction PPT for banking clients."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_GOLD = RGBColor(0xCC, 0x78, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
SECTION_BG = RGBColor(0x0A, 0x1A, 0x3A)
GREEN_ACCENT = RGBColor(0x2E, 0xCC, 0x71)
SUBTITLE_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
TERMINAL_BG = RGBColor(0x1E, 0x1E, 0x2E)
TERMINAL_GREEN = RGBColor(0x50, 0xFA, 0x7B)
TERMINAL_YELLOW = RGBColor(0xF1, 0xFA, 0x8C)
TERMINAL_RED = RGBColor(0xFF, 0x55, 0x55)
TERMINAL_CYAN = RGBColor(0x8B, 0xE9, 0xFD)
TERMINAL_WHITE = RGBColor(0xF8, 0xF8, 0xF2)
TERMINAL_ORANGE = RGBColor(0xFF, 0xB8, 0x6C)

PRODUCT_NAME = "基于大模型的数据库运维智能化平台"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

MONO_FONT = "Courier New"


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18,
             color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, spacing=Pt(8), bullet_char="\u25b8"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
    return txBox


def add_terminal_box(slide, left, top, width, height, lines, title="",
                     font_size=9, line_spacing=Pt(2)):
    """Add a dark terminal-style text box with colored lines.
    lines: list of (text, color) tuples or just strings (default white).
    """
    # Background
    bg_shape = add_rounded_rect(slide, left, top, width, height, TERMINAL_BG,
                                 line_color=RGBColor(0x44, 0x44, 0x66))

    # Title bar
    if title:
        add_shape(slide, left, top, width, Inches(0.32), RGBColor(0x2A, 0x2A, 0x45))
        add_text(slide, left + Inches(0.15), top + Inches(0.02), width - Inches(0.3),
                 Inches(0.28), title, font_size=10, color=TERMINAL_CYAN,
                 bold=True, font_name=MONO_FONT)
        content_top = top + Inches(0.35)
        content_height = height - Inches(0.4)
    else:
        content_top = top + Inches(0.08)
        content_height = height - Inches(0.16)

    txBox = slide.shapes.add_textbox(left + Inches(0.15), content_top,
                                      width - Inches(0.3), content_height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_item in enumerate(lines):
        if isinstance(line_item, tuple):
            text, color = line_item
        else:
            text, color = line_item, TERMINAL_WHITE

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = MONO_FONT
        p.space_after = line_spacing

    return txBox


def add_section_divider(text, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, SECTION_BG)
    add_shape(slide, Inches(5.5), Inches(3.0), Inches(2.3), Pt(3), ACCENT_GOLD)
    add_text(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(1.2),
             text, font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
                 subtitle, font_size=18, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


# ============================================================
# SLIDE 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GOLD)
add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
         PRODUCT_NAME, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
         "\u9762\u5411\u91d1\u878d\u7ea7\u751f\u4ea7\u73af\u5883\u7684\u6570\u636e\u5e93\u667a\u80fd\u8fd0\u7ef4\u89e3\u51b3\u65b9\u6848",
         font_size=22, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.5), Inches(4.5), Inches(2.3), Pt(2), ACCENT_GOLD)
tags = "\u672c\u5730\u90e8\u7f72  |  \u591a\u6a21\u578b\u9002\u914d  |  \u591a\u6570\u636e\u5e93\u8986\u76d6  |  \u65e0\u9700\u516c\u7f51  |  \u96f6\u4f9d\u8d56"
add_text(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
         tags, font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Pain Points & Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "\u91d1\u878d\u6570\u636e\u5e93\u8fd0\u7ef4\uff1a\u6311\u6218\u4e0e\u7834\u5c40",
         font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
         "\u5f53\u524d\u8fd0\u7ef4\u75db\u70b9", font_size=22, color=RGBColor(0xCC, 0x33, 0x33), bold=True)
pain_points = [
    "\u6570\u636e\u5e93\u89c4\u6a21\u6301\u7eed\u6269\u5f20\uff0cDBA \u4eba\u5747\u7ba1\u7406\u5b9e\u4f8b\u6570\u4e0d\u65ad\u6500\u5347",
    "\u6545\u969c\u8bca\u65ad\u4f9d\u8d56\u4e13\u5bb6\u7ecf\u9a8c\uff0c\u6392\u67e5\u5468\u671f\u957f\u3001\u7ed3\u8bba\u96be\u590d\u73b0",
    "\u76d1\u63a7\u544a\u8b66\u591a\u4f46\u65e0\u6839\u56e0\uff0c\u8fd0\u7ef4\u9677\u5165\"\u544a\u8b66\u75b2\u52b3\"",
    "\u591a\u5e93\u5f02\u6784\uff08Oracle/MySQL/PG\uff09\uff0c\u8fd0\u7ef4\u5de5\u5177\u788e\u7247\u5316",
    "\u91d1\u878d\u76d1\u7ba1\u8981\u6c42\uff1a\u6570\u636e\u4e0d\u51fa\u57df\uff0c\u5de5\u5177\u9700\u672c\u5730\u5316\u90e8\u7f72",
    "\u4f20\u7edf\u811a\u672c/\u5de5\u5177\u7f3a\u4e4f\u4e0a\u4e0b\u6587\u7406\u89e3\uff0c\u65e0\u6cd5\u6301\u7eed\u5b66\u4e60",
]
add_bullet_text(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.0),
                pain_points, font_size=15, color=DARK_TEXT, bullet_char="\u2715")

add_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
         "\u6211\u4eec\u7684\u89e3\u51b3\u65b9\u6848", font_size=22, color=GREEN_ACCENT, bold=True)
solutions = [
    "\u5927\u6a21\u578b\u9a71\u52a8\u7684\u591a\u8f6e\u94fe\u5f0f\u63a8\u7406\uff0c\u4ece\u544a\u8b66\u5230\u6839\u56e0\u5168\u81ea\u52a8\u95ed\u73af",
    "Sentinel \u5b9e\u65f6\u76d1\u6d4b + \u5f02\u5e38\u7a81\u53d1\u81ea\u52a8\u8bca\u65ad",
    "265+ \u89c4\u5219\u5f15\u64ce\u515c\u5e95\uff0c\u65e0\u5927\u6a21\u578b\u4e5f\u53ef\u8fd0\u884c",
    "\u56db\u5e93\u7edf\u4e00\u4ea4\u4e92\uff0c\u4e00\u5957\u5de5\u5177\u8986\u76d6 Oracle/MySQL/PG/OG",
    "\u5355\u4e8c\u8fdb\u5236\u672c\u5730\u90e8\u7f72\uff0c\u96f6\u5916\u90e8\u4f9d\u8d56\uff0c\u6570\u636e\u4e0d\u51fa\u751f\u4ea7\u7f51\u7edc",
    "\u4e0a\u4e0b\u6587\u8bb0\u5fc6\u7ba1\u7406\uff0c\u8d8a\u7528\u8d8a\u61c2\u4f60\u7684\u6570\u636e\u5e93",
]
add_bullet_text(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.0),
                solutions, font_size=15, color=DARK_TEXT, bullet_char="\u2713")

# ============================================================
# SLIDE 3: Architecture Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "\u5e73\u53f0\u67b6\u6784\u603b\u89c8", font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

layers = [
    ("\u4ea4\u4e92\u5c42", "\u547d\u4ee4\u884c /\u6307\u4ee4  |  SQL \u76f4\u67e5  |  \u81ea\u7136\u8bed\u8a00\u5bf9\u8bdd", ACCENT_BLUE, Inches(1.6)),
    ("\u667a\u80fd\u5f15\u64ce\u5c42", "\u591a\u8f6e\u94fe\u5f0f\u63a8\u7406\u5f15\u64ce  |  \u4e0a\u4e0b\u6587\u538b\u7f29  |  \u6d41\u5f0f\u8f93\u51fa  |  \u8bb0\u5fc6\u7ba1\u7406", RGBColor(0x2E, 0x86, 0xAB), Inches(2.8)),
    ("\u80fd\u529b\u5c42", "56+ \u5185\u7f6e\u5de5\u5177/\u6280\u80fd  |  265+ \u8bca\u65ad\u89c4\u5219  |  Sentinel \u5b9e\u65f6\u76d1\u6d4b  |  SQL Advisor", RGBColor(0x3B, 0x7D, 0xAD), Inches(4.0)),
    ("\u5927\u6a21\u578b\u9002\u914d\u5c42", "OpenAI \u517c\u5bb9\u534f\u8bae  |  Ollama \u672c\u5730\u6a21\u578b  |  \u591a\u5382\u5546\u9002\u914d\uff08\u901a\u4e49/DeepSeek/Kimi/...\uff09", RGBColor(0x48, 0x73, 0x9E), Inches(5.2)),
    ("\u6570\u636e\u5e93\u5c42", "Oracle  |  MySQL  |  PostgreSQL  |  OpenGauss", RGBColor(0x55, 0x69, 0x8F), Inches(6.4)),
]

for title, desc, color, top in layers:
    add_rounded_rect(slide, Inches(1.5), top, Inches(10.3), Inches(0.95), color)
    add_text(slide, Inches(1.8), top + Inches(0.08), Inches(2.5), Inches(0.4),
             title, font_size=18, color=WHITE, bold=True)
    add_text(slide, Inches(1.8), top + Inches(0.45), Inches(9.5), Inches(0.4),
             desc, font_size=13, color=RGBColor(0xDD, 0xDD, 0xDD))

for i in range(len(layers) - 1):
    add_text(slide, Inches(6.3), layers[i][3] + Inches(0.95), Inches(0.8), Inches(0.35),
             "\u25bc", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Three Interaction Modes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "\u4e09\u79cd\u4ea4\u4e92\u65b9\u5f0f\uff0c\u9002\u914d\u4e0d\u540c\u573a\u666f",
         font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

modes = [
    {
        "icon": "/",
        "title": "/\u6307\u4ee4\u6a21\u5f0f",
        "color": ACCENT_BLUE,
        "items": [
            "50+ \u5185\u7f6e\u547d\u4ee4\uff0c\u5feb\u901f\u7cbe\u51c6\u6267\u884c",
            "/health \u2014 \u4e00\u952e\u5168\u9762\u5065\u5eb7\u68c0\u67e5",
            "/dbtop \u2014 \u5b9e\u65f6\u6027\u80fd\u4eea\u8868\u76d8",
            "/sessions \u2014 \u4f1a\u8bdd\u5168\u666f\u89c6\u56fe",
            "/blocktree \u2014 \u9501\u94fe\u8def\u8ffd\u8e2a",
            "/slowsql \u2014 \u6162 SQL \u6355\u6349\u5206\u6790",
            "\u9002\u5408\uff1a\u65e5\u5e38\u5de1\u68c0\u3001\u5feb\u901f\u6392\u67e5",
        ]
    },
    {
        "icon": "SQL",
        "title": "SQL \u76f4\u67e5\u6a21\u5f0f",
        "color": RGBColor(0x2E, 0xCC, 0x71),
        "items": [
            "\u76f4\u63a5\u8f93\u5165 SQL \u8bed\u53e5\u6267\u884c",
            "\u5b89\u5168\u5206\u7ea7\u7ba1\u63a7\uff08\u53ea\u8bfb/\u64cd\u4f5c/\u7ba1\u7406/\u5371\u9669\uff09",
            "\u7ed3\u679c\u81ea\u52a8\u683c\u5f0f\u5316\u8868\u683c\u5c55\u793a",
            "\u652f\u6301 JSON / CSV \u683c\u5f0f\u5bfc\u51fa",
            "\u6267\u884c\u8ba1\u5212\u81ea\u52a8\u89e3\u6790\uff08/explain\uff09",
            "SQL Advisor \u4f18\u5316\u5efa\u8bae",
            "\u9002\u5408\uff1a\u6df1\u5ea6\u5206\u6790\u3001\u81ea\u5b9a\u4e49\u67e5\u8be2",
        ]
    },
    {
        "icon": "AI",
        "title": "\u81ea\u7136\u8bed\u8a00\u6a21\u5f0f",
        "color": ACCENT_GOLD,
        "items": [
            "\u7528\u81ea\u7136\u8bed\u8a00\u63cf\u8ff0\u95ee\u9898\u5373\u53ef",
            "\u5927\u6a21\u578b\u7406\u89e3\u610f\u56fe\u5e76\u7f16\u6392\u5de5\u5177",
            "\u4e09\u79cd\u63a8\u7406\u6df1\u5ea6\uff1a\u5feb\u8bca/\u8f85\u52a9/\u81ea\u52a8",
            "\u81ea\u52a8\u8c03\u7528\u5185\u7f6e\u5de5\u5177\u91c7\u96c6\u8bc1\u636e",
            "\u591a\u8f6e\u94fe\u5f0f\u63a8\u7406\u76f4\u8fbe\u6839\u56e0",
            "\u8f93\u51fa\u660e\u786e\u65b9\u6848 + \u53ef\u6267\u884c SQL",
            "\u9002\u5408\uff1a\u590d\u6742\u8bca\u65ad\u3001\u6839\u56e0\u5206\u6790",
        ]
    },
]

for i, mode in enumerate(modes):
    left = Inches(0.8 + i * 4.1)
    add_rounded_rect(slide, left, Inches(1.6), Inches(3.8), Inches(5.5),
                     CARD_BG, line_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_rounded_rect(slide, left, Inches(1.6), Inches(3.8), Inches(0.7), mode["color"])
    add_text(slide, left, Inches(1.65), Inches(3.8), Inches(0.6),
             f"{mode['icon']}  {mode['title']}", font_size=20,
             color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_text(slide, left + Inches(0.3), Inches(2.5), Inches(3.2), Inches(4.5),
                    mode["items"], font_size=13, color=DARK_TEXT, bullet_char="\u25b8")


# ============================================================
# SLIDE 5: Section - Real Cases
# ============================================================
add_section_divider("\u771f\u5b9e\u6848\u4f8b\u6f14\u793a",
                    "\u4ee5\u4e0b\u5747\u4e3a Oracle 19c \u751f\u4ea7\u73af\u5883\u5b9e\u6d4b\u7ed3\u679c")

# ============================================================
# SLIDE 6: Case 1 - /health Real Output
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
         "\u6848\u4f8b\u4e00\uff1a/health \u4e00\u952e\u5168\u9762\u5065\u5eb7\u68c0\u67e5",
         font_size=30, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(0.88), Inches(1.7), Pt(3), ACCENT_GOLD)

# Tag
add_text(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.35),
         "\u73af\u5883\uff1aOracle 19c (ORCLCDB)  |  \u4e00\u6761\u547d\u4ee4\uff0c20+ \u68c0\u67e5\u9879\uff0c7 \u4e2a\u7ef4\u5ea6\u5168\u8986\u76d6",
         font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Terminal box - Real /health output
health_lines = [
    ("$ /health", TERMINAL_GREEN),
    ("", TERMINAL_WHITE),
    ("\u250c\u2500 Database Health \u2500 ORCLCDB \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510", TERMINAL_CYAN),
    ("\u2502  Overall: \u2717 CRITICAL  (6 issues found)                                   \u2502", TERMINAL_RED),
    ("\u251c\u2500 Basic \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524", TERMINAL_CYAN),
    ("\u2502  Instance       UP 2.4 hours       \u2713   DB Role       PRIMARY          \u2713 \u2502", TERMINAL_GREEN),
    ("\u2502  Archive Mode   NOARCHIVELOG       \u26a0                                     \u2502", TERMINAL_YELLOW),
    ("\u251c\u2500 Storage \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524", TERMINAL_CYAN),
    ("\u2502  Tablespace     Max SYSTEM 2.2%    \u2713   Temp Space    Max 0.0%         \u2713 \u2502", TERMINAL_GREEN),
    ("\u2502  Undo Space     No data            \u2713   FRA Usage     Not configured   \u2713 \u2502", TERMINAL_GREEN),
    ("\u251c\u2500 Session/Perf \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524", TERMINAL_CYAN),
    ("\u2502  Connections    83 / 2000 (4%)     \u2713   Active Sess   2                \u2713 \u2502", TERMINAL_GREEN),
    ("\u2502  Wait Events    Top1: row cache l\u2026 \u2717   Slow SQL      1 avg >5s       \u26a0 \u2502", TERMINAL_RED),
    ("\u2502  Buffer Hit     100.0%             \u2713   Library Hit   99.8%            \u2713 \u2502", TERMINAL_GREEN),
    ("\u251c\u2500 Memory \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524", TERMINAL_CYAN),
    ("\u2502  PGA Usage      7 MB (auto)        \u2713   SGA Overview  9216 MB          \u2713 \u2502", TERMINAL_GREEN),
    ("\u2502  Shared Pool    Free 61.8%         \u2713                                     \u2502", TERMINAL_GREEN),
    ("\u251c\u2500 Logs \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524", TERMINAL_CYAN),
    ("\u2502  Redo Switch    17x in last 1h     \u26a0   Alert Log     No alerts 24h   \u2713 \u2502", TERMINAL_YELLOW),
    ("\u251c\u2500 Maintenance \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524", TERMINAL_CYAN),
    ("\u2502  Backup         No backup records  \u26a0   Statistics    3.7% stale      \u2713 \u2502", TERMINAL_YELLOW),
    ("\u2502  Invalid Obj    1 object           \u26a0   Resource Lim  No limits hit   \u2713 \u2502", TERMINAL_YELLOW),
    ("\u2502  Password Exp   None expiring      \u2713                                     \u2502", TERMINAL_GREEN),
    ("\u251c\u2500 Alerts \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524", TERMINAL_CYAN),
    ("\u2502  \u26a0 Archive mode: NOARCHIVELOG                                            \u2502", TERMINAL_YELLOW),
    ("\u2502  \u2717 Wait event: Top1: row cache lock 2653.5s                              \u2502", TERMINAL_RED),
    ("\u2502  \u26a0 Slow SQL: 1 avg >5s                                                   \u2502", TERMINAL_YELLOW),
    ("\u2502  Tip: /waits to inspect; /slowsql to find slow SQL; /redo to check      \u2502", TERMINAL_WHITE),
    ("\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518", TERMINAL_CYAN),
]
add_terminal_box(slide, Inches(0.8), Inches(1.45), Inches(12.0), Inches(5.8),
                 health_lines, title="Terminal \u2014 Oracle 19c ORCLCDB", font_size=9)

# ============================================================
# SLIDE 7: Case 2 - /dbtop Real Output
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
         "\u6848\u4f8b\u4e8c\uff1a/dbtop \u5b9e\u65f6\u6027\u80fd\u4eea\u8868\u76d8\uff08108 \u6d3b\u8dc3\u4f1a\u8bdd\u538b\u6d4b\uff09",
         font_size=30, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(0.88), Inches(1.7), Pt(3), ACCENT_GOLD)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.35),
         "\u73af\u5883\uff1aOracle 19c  |  \u538b\u6d4b\u4e2d 108 \u6d3b\u8dc3\u4f1a\u8bdd  |  \u7c7b\u4f3c Linux top \u7684\u6570\u636e\u5e93\u5b9e\u65f6\u76d1\u63a7",
         font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

dbtop_lines = [
    ("$ /dbtop", TERMINAL_GREEN),
    ("", TERMINAL_WHITE),
    ("\u256d\u2500 dbtop \u2500\u2500 oracle 19c \u2500\u2500 ORCLCDB \u2500\u2500 PRIMARY \u2500\u2500 \u25cf WARNING \u2500\u2500 15:42:07 \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u256e", TERMINAL_CYAN),
    ("\u2502                                                                            \u2502", TERMINAL_CYAN),
    ("\u2502  SGA \u2587\u2587\u2587\u2587\u2587\u2587\u2591\u2591 6,834M   PGA \u2587\u2587\u2591\u2591\u2591\u2591\u2591\u2591 247M   db% \u2587\u2587\u2587\u2587\u2587\u2587\u2587\u2591 78.3   WTR% \u2587\u2587\u2587\u2587\u2587\u2587\u2587\u2587 92.1  \u2502", TERMINAL_WHITE),
    ("\u2502  Session 212  Active 108  ActiveCPU 6  ActiveIO 3  |  TPS 1,247  QPS 8,934 \u2502", TERMINAL_WHITE),
    ("\u2570\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u256f", TERMINAL_CYAN),
    ("\u256d\u2500 Top Wait Events \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u256e", TERMINAL_CYAN),
    ("\u2502  EVENT                    Dwait Dtime(ms) DPCT   Delta                      \u2502", TERMINAL_WHITE),
    ("\u2502  enq: SQ - contention     1,247  35,812  \u2587\u2587\u2587\u2587\u2587\u2587 46.1%                      \u2502", TERMINAL_RED),
    ("\u2502  buffer busy waits         2,891  22,658  \u2587\u2587\u2587\u2587\u2591\u2591 36.7%                      \u2502", TERMINAL_ORANGE),
    ("\u2502  log file switch (chec..     12   6,607  \u2587\u2587\u2591\u2591\u2591\u2591 10.7%                      \u2502", TERMINAL_YELLOW),
    ("\u2502  cursor: pin S               892  1,914  \u2587\u2591\u2591\u2591\u2591\u2591  3.1%                      \u2502", TERMINAL_WHITE),
    ("\u2570\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u256f", TERMINAL_CYAN),
    ("\u256d\u2500 Active Sessions (108) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u256e", TERMINAL_CYAN),
    ("\u2502  SID   USR        SQLID         EVENT                CLASS       E/T SQL  \u2502", TERMINAL_WHITE),
    ("\u2502  11    LOADTEST   4j031tn1y64s2 enq: SQ-contention  Config     3.2s INS..\u2502", TERMINAL_RED),
    ("\u2502  12    LOADTEST   4j031tn1y64s2 enq: SQ-contention  Config     2.8s INS..\u2502", TERMINAL_RED),
    ("\u2502  15    LOADTEST   fbnsszmxf8zsj cursor: pin S        Concur     1.5s BEGI.\u2502", TERMINAL_YELLOW),
    ("\u2502  17    LOADTEST   4s20v2f39mp83 buffer busy waits    Concur     0.8s upd..\u2502", TERMINAL_ORANGE),
    ("\u2502  20    LOADTEST   4j031tn1y64s2 enq: SQ-contention  Config     4.1s INS..\u2502", TERMINAL_RED),
    ("\u2502  ... and 103 more active sessions                                          \u2502", TERMINAL_WHITE),
    ("\u2570\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u256f", TERMINAL_CYAN),
]
add_terminal_box(slide, Inches(0.6), Inches(1.45), Inches(12.2), Inches(5.8),
                 dbtop_lines, title="Terminal \u2014 /dbtop \u5b9e\u65f6\u76d1\u63a7", font_size=9)


# ============================================================
# SLIDE 8: Case 3 - Sentinel Anomaly Detection (Real E2E)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
         "\u6848\u4f8b\u4e09\uff1aSentinel \u5b9e\u65f6\u5f02\u5e38\u68c0\u6d4b\uff08\u771f\u5b9e\u538b\u6d4b\u573a\u666f\uff09",
         font_size=30, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(0.88), Inches(1.7), Pt(3), ACCENT_GOLD)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.35),
         "\u73af\u5883\uff1aOracle 19c  |  \u538b\u6d4b\u6d41\u91cf: 16 \u5e76\u53d1 workers (heavy/io/dml/lock \u5404 4)  |  \u6d3b\u8dc3\u4f1a\u8bdd\u4ece 8 \u7a81\u589e\u5230 41",
         font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

sentinel_lines = [
    ("\u25b8 Active Sessions 8.0\u219241 (3\u03c3 threshold 24.0)", TERMINAL_RED),
    ("", TERMINAL_WHITE),
    ("\u2500\u2500 Report #1 \u2500\u2500", TERMINAL_CYAN),
    ("  Trigger: active_sessions  8.0 \u2192 41.0", TERMINAL_WHITE),
    ("  Duration: 39.8s", TERMINAL_WHITE),
    ("", TERMINAL_WHITE),
    ("  Wait Event Distribution:", TERMINAL_CYAN),
    ("  \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510", TERMINAL_CYAN),
    ("  \u2502 EVENT                    \u2502 PCT   \u2502 Distribution       \u2502 WAIT_CLASS   \u2502", TERMINAL_WHITE),
    ("  \u2502 enq: SQ - contention     \u2502 53.6% \u2502 \u2588\u2588\u2588\u2588\u2588\u2588\u2588\u2588\u2588\u2588\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591 \u2502 Configuration\u2502", TERMINAL_RED),
    ("  \u2502 DB CPU                   \u2502 27.0% \u2502 \u2588\u2588\u2588\u2588\u2588\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591 \u2502 CPU          \u2502", TERMINAL_ORANGE),
    ("  \u2502 buffer busy waits        \u2502 13.9% \u2502 \u2588\u2588\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591\u2591 \u2502 Concurrency  \u2502", TERMINAL_YELLOW),
    ("  \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518", TERMINAL_CYAN),
    ("", TERMINAL_WHITE),
    ("  Top SQL:", TERMINAL_CYAN),
    ("  \u250c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510", TERMINAL_CYAN),
    ("  \u2502 SQL_ID        \u2502 Elapsed \u2502 Conc\u2502 Wait Event                \u2502", TERMINAL_WHITE),
    ("  \u2502 4j031tn1y64s2 \u2502 1.0s    \u2502 40  \u2502 enq: SQ - contention      \u2502", TERMINAL_RED),
    ("  \u2502 4s20v2f39mp83 \u2502 1.0s    \u2502 18  \u2502 buffer busy waits         \u2502", TERMINAL_ORANGE),
    ("  \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518", TERMINAL_CYAN),
    ("", TERMINAL_WHITE),
    ("  Blocking Chains:", TERMINAL_CYAN),
    ("  \u2502 SID 2275 \u2502 SQL: 4j031tn1y64s2 \u2502 enq: SQ - contention \u2502 Victims: 39 \u2502", TERMINAL_RED),
    ("  \u2502 SID 1534 \u2502 SQL: 4j031tn1y64s2 \u2502 enq: SQ - contention \u2502 Victims: 37 \u2502", TERMINAL_RED),
    ("  \u2502 SID 2663 \u2502 SQL: 4s20v2f39mp83 \u2502 buffer busy waits    \u2502 Victims: 29 \u2502", TERMINAL_ORANGE),
]
add_terminal_box(slide, Inches(0.6), Inches(1.45), Inches(12.2), Inches(5.8),
                 sentinel_lines, title="Terminal \u2014 Sentinel \u5f02\u5e38\u68c0\u6d4b\u62a5\u544a", font_size=9)


# ============================================================
# SLIDE 9: Case 4 - /llm AI Deep Diagnosis (Real E2E)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
         "\u6848\u4f8b\u56db\uff1a/llm \u5927\u6a21\u578b\u591a\u8f6e\u94fe\u5f0f\u63a8\u7406\u8bca\u65ad\uff08\u771f\u5b9e\u538b\u6d4b\uff09",
         font_size=30, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(0.88), Inches(1.7), Pt(3), ACCENT_GOLD)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.35),
         "\u6a21\u578b\uff1aClaude Opus  |  3 \u8f6e\u63a8\u7406  |  \u8017\u65f6 1 \u5206 59 \u79d2  |  \u81ea\u52a8\u8c03\u7528 /health\u3001/waits \u91c7\u96c6\u8bc1\u636e",
         font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

llm_lines = [
    ("$ /llm", TERMINAL_GREEN),
    ("", TERMINAL_WHITE),
    ("  \u273b Done (1m 59s)", TERMINAL_GREEN),
    ("  \u251c Round 1: called health (21.7s)", TERMINAL_CYAN),
    ("  \u251c Round 2: called waits (1m 5s)", TERMINAL_CYAN),
    ("  \u251c Round 3: output final diagnosis (32.0s)", TERMINAL_CYAN),
    ("", TERMINAL_WHITE),
    ("\u25a0 Database Status Diagnosis", TERMINAL_CYAN),
    ("", TERMINAL_WHITE),
    ("  \u25b8 Issue 1 (Critical): Sequence Contention \u2014 70.6% of total waits", TERMINAL_RED),
    ("  \u250c\u2500\u2500\u2500\u2500\u2500\u252c\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510", TERMINAL_CYAN),
    ("  \u2502 Root\u2502 ORDER_ITEMS.ITEM_ID identity column sequence cache only 20,    \u2502", TERMINAL_WHITE),
    ("  \u2502 Cause\u2502 243.8K contentions, cumulative wait 3630.7s                    \u2502", TERMINAL_WHITE),
    ("  \u2502 Evid\u2502 51 sessions waiting on enq: SQ-contention,                     \u2502", TERMINAL_WHITE),
    ("  \u2502     \u2502 SQL g11yt3q86jfx2 (INSERT INTO ORDER_ITEMS) avg exec 5.18s     \u2502", TERMINAL_WHITE),
    ("  \u2514\u2500\u2500\u2500\u2500\u2500\u2534\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518", TERMINAL_CYAN),
    ("", TERMINAL_WHITE),
    ("  Immediate fix \u2014 increase sequence cache:", TERMINAL_GREEN),
    ("  \u250c\u2500 sql \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510", TERMINAL_CYAN),
    ("  \u2502 ALTER TABLE LOADUSER.ORDER_ITEMS MODIFY ITEM_ID                          \u2502", TERMINAL_YELLOW),
    ("  \u2502   GENERATED ALWAYS AS IDENTITY (CACHE 5000 NOORDER);                    \u2502", TERMINAL_YELLOW),
    ("  \u2502 ALTER TABLE LOADUSER.ORDERS MODIFY ORDER_ID                              \u2502", TERMINAL_YELLOW),
    ("  \u2502   GENERATED ALWAYS AS IDENTITY (CACHE 5000 NOORDER);                    \u2502", TERMINAL_YELLOW),
    ("  \u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518", TERMINAL_CYAN),
    ("", TERMINAL_WHITE),
    ("  \u25b8 Issue 2 (Important): Buffer Busy Waits \u2014 21.4% of total waits", TERMINAL_ORANGE),
    ("  \u25b8 Issue 3 (Critical): Redo Log Switches Too Frequent (22 sec/switch)", TERMINAL_ORANGE),
    ("", TERMINAL_WHITE),
    ("  \u25b8 Execution Priority:", TERMINAL_GREEN),
    ("  1. Immediate: Increase identity column cache \u2192 resolve 51 session blocks", TERMINAL_WHITE),
    ("  2. Immediate: Expand redo logs to 1GB \u2192 eliminate frequent switches", TERMINAL_WHITE),
    ("  3. Observe: buffer busy waits \u2014 retest after first two fixes", TERMINAL_WHITE),
]
add_terminal_box(slide, Inches(0.6), Inches(1.45), Inches(12.2), Inches(5.8),
                 llm_lines, title="Terminal \u2014 /llm \u5927\u6a21\u578b\u8bca\u65ad (Claude Opus, 3 \u8f6e\u63a8\u7406)", font_size=9)


# ============================================================
# SLIDE 10: Case 5 - /rule Fast Diagnosis (No LLM)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
         "\u6848\u4f8b\u4e94\uff1a/rule \u89c4\u5219\u5f15\u64ce\u5feb\u901f\u8bca\u65ad\uff08\u65e0\u9700\u5927\u6a21\u578b\uff0c2.7 \u79d2\u5b8c\u6210\uff09",
         font_size=30, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(0.88), Inches(1.7), Pt(3), ACCENT_GOLD)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.35),
         "265+ \u89c4\u5219  |  \u7eaf\u672c\u5730\u6267\u884c  |  \u65e0\u7f51\u7edc\u4f9d\u8d56  |  \u5927\u6a21\u578b\u4e0d\u53ef\u7528\u65f6\u81ea\u52a8\u964d\u7ea7",
         font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

rule_lines = [
    ("$ /rule", TERMINAL_GREEN),
    ("", TERMINAL_WHITE),
    ("\u2500\u2500 Rule Diagnosis \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", TERMINAL_CYAN),
    ("", TERMINAL_WHITE),
    ("  Root Cause: Sequence CACHE too small or NOCACHE causing SQ contention", TERMINAL_RED),
    ("  Severity: \u25a0\u25a0\u25a0\u25a1 HIGH    Confidence: 78%", TERMINAL_ORANGE),
    ("", TERMINAL_WHITE),
    ("  \u2500\u2500 Evidence Chain \u2500\u2500", TERMINAL_CYAN),
    ("  1. enq: SQ - contention in single-instance: small sequence CACHE", TERMINAL_WHITE),
    ("  2. Each CACHE exhaustion updates data dictionary seq$,", TERMINAL_WHITE),
    ("     causing row cache lock and SQ enqueue", TERMINAL_WHITE),
    ("  3. High-concurrency INSERT with NOCACHE sequences especially severe", TERMINAL_WHITE),
    ("  4. Related rule WE2-007 (enq: SQ) check recommended", TERMINAL_WHITE),
    ("  5. Related rule WD015 (row cache lock) check recommended", TERMINAL_WHITE),
    ("", TERMINAL_WHITE),
    ("  \u2500\u2500 Remediation \u2500\u2500", TERMINAL_GREEN),
    ("  [URGENT] 1. Find contended sequences and increase CACHE", TERMINAL_YELLOW),
    ("           > /sql \"SELECT s.sequence_owner, s.sequence_name,", TERMINAL_YELLOW),
    ("             s.cache_size, s.order_flag, s.last_number", TERMINAL_YELLOW),
    ("             FROM dba_sequences s", TERMINAL_YELLOW),
    ("             WHERE s.sequence_owner NOT IN ('SYS','SYSTEM')\"", TERMINAL_YELLOW),
    ("           Risk: sequence number gaps increase (uniqueness unaffected)", TERMINAL_WHITE),
    ("", TERMINAL_WHITE),
    ("  [FIX]    2. Check ASH to confirm contended sequence objects", TERMINAL_YELLOW),
    ("           > /sql \"SELECT h.current_obj#, o.object_name, COUNT(*)", TERMINAL_YELLOW),
    ("             FROM v$active_session_history h LEFT JOIN dba_objects o", TERMINAL_YELLOW),
    ("             ON h.current_obj# = o.object_id", TERMINAL_YELLOW),
    ("             WHERE h.event = 'enq: SQ - contention'\"", TERMINAL_YELLOW),
    ("", TERMINAL_WHITE),
    ("  \u2500\u2500 Correlated Analysis \u2500\u2500", TERMINAL_CYAN),
    ("  \u2022 cursor: pin S: downstream symptom, resolves after fixing root cause", TERMINAL_WHITE),
    ("  \u2022 buffer busy waits: downstream symptom, resolves after fixing root cause", TERMINAL_WHITE),
    ("", TERMINAL_WHITE),
    ("\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500", TERMINAL_CYAN),
    ("  Completed in 2.7s (265 rules evaluated)", TERMINAL_GREEN),
]
add_terminal_box(slide, Inches(0.6), Inches(1.45), Inches(12.2), Inches(5.8),
                 rule_lines, title="Terminal \u2014 /rule \u89c4\u5219\u5f15\u64ce\u8bca\u65ad (2.7s, \u65e0\u9700 LLM)", font_size=9)


# ============================================================
# SLIDE 11: Case 6 - E2E: Sentinel → LLM Diagnosis (Lock Contention)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.6),
         "\u6848\u4f8b\u516d\uff1aSentinel \u5f02\u5e38\u68c0\u6d4b \u2192 \u5927\u6a21\u578b\u81ea\u52a8\u8bca\u65ad\uff08\u5168\u94fe\u8def\uff09",
         font_size=30, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(0.88), Inches(1.7), Pt(3), ACCENT_GOLD)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.35),
         "\u573a\u666f\uff1a16 \u5e76\u53d1\u7ebf\u7a0b\u4ea7\u751f\u884c\u9501\u7ade\u4e89  |  \u6a21\u578b: Qwen 3.5 9B (\u672c\u5730 Ollama)  |  \u8017\u65f6 115s  |  Playbook \u6a21\u5f0f",
         font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Left side: Detection flow
add_text(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4),
         "\u2460 Sentinel \u5f02\u5e38\u68c0\u6d4b\u4e0e\u8bc1\u636e\u91c7\u96c6", font_size=18, color=ACCENT_BLUE, bold=True)

e2e_detect_lines = [
    ("\u25b6 \u57fa\u7ebf\u91c7\u96c6 (6 \u6837\u672c, 1s \u95f4\u9694)", TERMINAL_CYAN),
    ("  [1] active=12 db%=0.0    TPS=0    QPS=0", TERMINAL_WHITE),
    ("  [2] active=11 db%=1324.5 TPS=807  QPS=32338", TERMINAL_WHITE),
    ("  [6] active=14 db%=1130.6 TPS=3    QPS=2077", TERMINAL_WHITE),
    ("", TERMINAL_WHITE),
    ("\u25b6 \u5f02\u5e38\u68c0\u6d4b", TERMINAL_CYAN),
    ("  \u5206\u7c7b: \u9501\u7b49\u5f85\u963b\u585e (lock_contention)", TERMINAL_RED),
    ("  \u7f6e\u4fe1\u5ea6: 70%  \u5cf0\u503c: 15  \u57fa\u7ebf: 12.3", TERMINAL_ORANGE),
    ("", TERMINAL_WHITE),
    ("\u25b6 \u8bc1\u636e (\u963b\u585e\u94fe)", TERMINAL_CYAN),
    ("  SID 1002 \u963b\u585e 3 \u4e2a\u4f1a\u8bdd, enq: TX-row lock", TERMINAL_RED),
    ("  SID 1142 \u963b\u585e 3 \u4e2a\u4f1a\u8bdd, enq: TX-row lock", TERMINAL_RED),
    ("", TERMINAL_WHITE),
    ("\u25b6 Top SQL", TERMINAL_CYAN),
    ("  UPDATE opendb_lock_test SET val=:1", TERMINAL_WHITE),
    ("  \u51fa\u73b0\u7387=75%  \u5e76\u53d1=3  \u8017\u65f6=4.0s", TERMINAL_WHITE),
    ("", TERMINAL_WHITE),
    ("\u25b6 \u7b49\u5f85\u5206\u5e03", TERMINAL_CYAN),
    ("  CPU: 74.0%  Application: 24.7%", TERMINAL_WHITE),
]
add_terminal_box(slide, Inches(0.6), Inches(1.95), Inches(5.8), Inches(5.2),
                 e2e_detect_lines, title="Sentinel \u81ea\u52a8\u68c0\u6d4b", font_size=9)

# Right side: LLM diagnosis
add_text(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.4),
         "\u2461 \u5927\u6a21\u578b\u63a8\u7406\u8f93\u51fa (Qwen 3.5 9B)", font_size=18, color=ACCENT_GOLD, bold=True)

e2e_diag_lines = [
    ("\u25a0 \u6839\u56e0\u5206\u6790", TERMINAL_CYAN),
    ("  \u6839\u672c\u539f\u56e0: TX\u884c\u9501\u7ade\u4e89 (lock_contention)", TERMINAL_RED),
    ("", TERMINAL_WHITE),
    ("  1. enq: TX-row lock contention \u662f\u5178\u578b", TERMINAL_WHITE),
    ("     \u884c\u7ea7\u9501\u7b49\u5f85\u4e8b\u4ef6", TERMINAL_WHITE),
    ("  2. UPDATE opendb_lock_test \u51fa\u73b0\u7387 75%\u3001", TERMINAL_WHITE),
    ("     \u8017\u65f6 4.0s\uff0c\u8be5\u8868\u662f\u9501\u7ade\u4e89\u4e3b\u8981\u6765\u6e90", TERMINAL_WHITE),
    ("  3. \u6d3b\u8dc3\u4f1a\u8bdd\u7a81\u589e\u8bf1\u56e0: \u963b\u585e\u5bfc\u81f4", TERMINAL_WHITE),
    ("     \u5176\u4ed6 SQL \u7b49\u5f85\u5728 enq: TX \u4e0a\u65e0\u6cd5\u7ed3\u675f", TERMINAL_WHITE),
    ("", TERMINAL_WHITE),
    ("\u25a0 \u7d27\u6025\u63aa\u65bd\uff08\u6b62\u8840\u65b9\u6848\uff09", TERMINAL_GREEN),
    ("  SELECT sid, serial#, sql_id, event", TERMINAL_YELLOW),
    ("  FROM v$session WHERE event LIKE", TERMINAL_YELLOW),
    ("  '%enq: TX%' AND blocking_session", TERMINAL_YELLOW),
    ("  IS NOT NULL;", TERMINAL_YELLOW),
    ("  ALTER SYSTEM KILL SESSION 'sid,s#';", TERMINAL_YELLOW),
    ("", TERMINAL_WHITE),
    ("\u25a0 \u6839\u56e0\u4fee\u590d\uff08\u957f\u671f\u65b9\u6848\uff09", TERMINAL_GREEN),
    ("  CREATE INDEX idx_opendb_lock_test_id", TERMINAL_YELLOW),
    ("    ON opendb_lock_test(id);", TERMINAL_YELLOW),
    ("  \u51cf\u5c11\u4e8b\u52a1\u6301\u6709\u9501\u65f6\u95f4\uff1a\u5206\u6279\u63d0\u4ea4", TERMINAL_WHITE),
    ("  \u9694\u79bb\u7ea7\u522b\u8c03\u6574: READ COMMITTED", TERMINAL_WHITE),
]
add_terminal_box(slide, Inches(6.8), Inches(1.95), Inches(5.8), Inches(5.2),
                 e2e_diag_lines, title="LLM \u8bca\u65ad\u8f93\u51fa (115s, 1405 tokens)", font_size=9)


# ============================================================
# SLIDE 12: Section - Feature Details
# ============================================================
add_section_divider("\u5e73\u53f0\u80fd\u529b\u8be6\u89e3",
                    "360\u00b0 \u5168\u65b9\u4f4d\u6570\u636e\u5e93\u5256\u6790 \u00b7 \u667a\u80fd\u8bca\u65ad \u00b7 \u8bb0\u5fc6\u7ba1\u7406")


# ============================================================
# SLIDE 13: 360-degree Coverage
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "360\u00b0 \u5168\u65b9\u4f4d\u6570\u636e\u5e93\u5256\u6790", font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

categories = [
    {
        "title": "\u6027\u80fd\u76d1\u63a7",
        "color": ACCENT_BLUE,
        "items": ["/dbtop \u5b9e\u65f6\u4eea\u8868\u76d8", "/awr \u8d1f\u8f7d\u62a5\u544a\u5206\u6790", "/ash \u6d3b\u8dc3\u4f1a\u8bdd\u91c7\u6837", "/perfsnap \u6027\u80fd\u5feb\u7167\u5bf9\u6bd4"]
    },
    {
        "title": "\u4f1a\u8bdd\u7ba1\u7406",
        "color": RGBColor(0x2E, 0x86, 0xAB),
        "items": ["/sessions \u5168\u666f\u89c6\u56fe", "/activesessions \u6d3b\u8dc3\u4f1a\u8bdd", "/kill \u4f1a\u8bdd\u7ec8\u6b62", "/blocktree \u9501\u94fe\u8def\u5206\u6790"]
    },
    {
        "title": "\u7b49\u5f85\u5206\u6790",
        "color": RGBColor(0x9B, 0x59, 0xB6),
        "items": ["/waits \u7b49\u5f85\u4e8b\u4ef6\u7edf\u8ba1", "/latches \u95e9\u9501\u4e89\u7528", "/mutexes \u4e92\u65a5\u4e89\u7528", "/locks \u884c/\u8868\u9501\u76d1\u63a7"]
    },
    {
        "title": "SQL \u5206\u6790",
        "color": RGBColor(0xE6, 0x7E, 0x22),
        "items": ["/slowsql \u6162 SQL \u6355\u6349", "/topsql \u9ad8\u8017 SQL \u6392\u884c", "/explain \u6267\u884c\u8ba1\u5212\u89e3\u6790", "/sqladvisor \u4f18\u5316\u5efa\u8bae"]
    },
    {
        "title": "\u5b58\u50a8\u4e0e\u5185\u5b58",
        "color": RGBColor(0x2E, 0xCC, 0x71),
        "items": ["/space \u8868\u7a7a\u95f4\u4f7f\u7528\u7387", "/sga SGA \u5185\u5b58\u8be6\u60c5", "/pga PGA \u5185\u5b58\u8be6\u60c5", "/segments \u5927\u5bf9\u8c61\u5206\u6790"]
    },
    {
        "title": "\u5065\u5eb7\u4e0e\u5b89\u5168",
        "color": RGBColor(0xCC, 0x33, 0x33),
        "items": ["/health \u5065\u5eb7\u68c0\u67e5", "/indexhealth \u7d22\u5f15\u5065\u5eb7", "/users \u8d26\u6237\u5230\u671f\u68c0\u67e5", "/resource \u8d44\u6e90\u9650\u5236"]
    },
    {
        "title": "\u65e5\u5fd7\u4e0e\u6062\u590d",
        "color": RGBColor(0x1A, 0xBC, 0x9C),
        "items": ["/redo \u91cd\u505a\u65e5\u5fd7\u72b6\u6001", "/fra \u5feb\u901f\u6062\u590d\u533a", "/alert \u544a\u8b66\u65e5\u5fd7", "/backup \u5907\u4efd\u5386\u53f2"]
    },
    {
        "title": "\u9ad8\u53ef\u7528",
        "color": RGBColor(0x34, 0x49, 0x5E),
        "items": ["/standby DG \u72b6\u6001", "/asm \u78c1\u76d8\u7ec4\u4fe1\u606f", "/trace \u5806\u6808\u8ffd\u8e2a", "/scheduler \u5b9a\u65f6\u4efb\u52a1"]
    },
]

for i, cat in enumerate(categories):
    col = i % 4
    row = i // 4
    left = Inches(0.6 + col * 3.15)
    top = Inches(1.5 + row * 2.9)
    add_rounded_rect(slide, left, top, Inches(2.9), Inches(2.6),
                     CARD_BG, line_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, left, top, Inches(2.9), Inches(0.06), cat["color"])
    add_text(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.5), Inches(0.4),
             cat["title"], font_size=16, color=cat["color"], bold=True)
    for j, item in enumerate(cat["items"]):
        add_text(slide, left + Inches(0.2), top + Inches(0.65 + j * 0.42), Inches(2.5), Inches(0.35),
                 f"\u25b8 {item}", font_size=12, color=DARK_TEXT)


# ============================================================
# SLIDE 14: Context & Memory Management
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "\u4e0a\u4e0b\u6587\u4e0e\u8bb0\u5fc6\u7ba1\u7406\uff1a\u8d8a\u7528\u8d8a\u61c2\u4f60\u7684\u6570\u636e\u5e93",
         font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

# 6-layer context
add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "\u516d\u5c42\u4e0a\u4e0b\u6587\u67b6\u6784", font_size=22, color=ACCENT_BLUE, bold=True)

context_layers = [
    ("\u901a\u7528\u7cfb\u7edf\u63d0\u793a", "\u57fa\u7840\u6307\u4ee4\u96c6\uff0c\u786e\u4fdd\u5b89\u5168\u53ef\u63a7", ACCENT_BLUE),
    ("\u4ea7\u54c1\u4e13\u5c5e\u89c4\u5219", "Oracle/MySQL/PG/OG \u7279\u5b9a\u77e5\u8bc6", RGBColor(0x2E, 0x86, 0xAB)),
    ("\u538b\u7f29\u8bca\u65ad\u62a5\u544a", "Sentinel \u7a81\u53d1\u4e8b\u4ef6\u6458\u8981", RGBColor(0x3B, 0x7D, 0xAD)),
    ("\u5b9e\u4f8b\u8bb0\u5fc6", "\u8be5\u6570\u636e\u5e93\u7684\u5386\u53f2\u753b\u50cf\u4e0e\u8bca\u65ad\u7ed3\u8bba", RGBColor(0x48, 0x73, 0x9E)),
    ("\u7b56\u7565\u5c42", "\u5168\u5c40\u2192\u7cfb\u7edf\u2192\u5b9e\u4f8b\u2192\u7528\u6237 \u56db\u7ea7\u7b56\u7565", RGBColor(0x55, 0x69, 0x8F)),
    ("\u5de5\u5177\u8c03\u7528\u5386\u53f2", "\u672c\u8f6e\u63a8\u7406\u7684\u8bc1\u636e\u94fe\u8ffd\u8e2a", RGBColor(0x62, 0x5F, 0x80)),
]

for i, (name, desc, color) in enumerate(context_layers):
    top = Inches(2.1 + i * 0.68)
    add_rounded_rect(slide, Inches(0.8), top, Inches(5.5), Inches(0.58), color)
    add_text(slide, Inches(1.0), top + Inches(0.03), Inches(2.0), Inches(0.25),
             name, font_size=14, color=WHITE, bold=True)
    add_text(slide, Inches(1.0), top + Inches(0.28), Inches(5.0), Inches(0.25),
             desc, font_size=11, color=RGBColor(0xDD, 0xDD, 0xDD))

add_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
         "\u5b9e\u4f8b\u8bb0\u5fc6\u7cfb\u7edf", font_size=22, color=ACCENT_BLUE, bold=True)

memory_features = [
    "\u6bcf\u4e2a\u6570\u636e\u5e93\u5b9e\u4f8b\u72ec\u7acb\u7684\u77e5\u8bc6\u5e93",
    "\u81ea\u52a8\u79ef\u7d2f\uff1a\u914d\u7f6e\u753b\u50cf\u3001\u67b6\u6784\u7279\u5f81\u3001\u5386\u53f2\u95ee\u9898",
    "\u8bca\u65ad\u7ed3\u8bba\u6c89\u6dc0\uff1a\u76f8\u4f3c\u95ee\u9898\u5feb\u901f\u5173\u8054\u5386\u53f2",
    "\u7528\u5f97\u8d8a\u4e45\uff0c\u5927\u6a21\u578b\u5bf9\u8be5\u5e93\u8d8a\u4e86\u89e3",
    "\u8de8\u4f1a\u8bdd\u6301\u4e45\u5316\uff0c\u91cd\u542f\u4e0d\u4e22\u5931",
    "\u652f\u6301\u591a\u5b9e\u4f8b\u7ba1\u7406\uff0c\u5207\u6362\u81ea\u52a8\u52a0\u8f7d\u5bf9\u5e94\u8bb0\u5fc6",
]
add_bullet_text(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.0),
                memory_features, font_size=14, color=DARK_TEXT)

add_rounded_rect(slide, Inches(7.0), Inches(4.8), Inches(5.5), Inches(2.3),
                 RGBColor(0xFE, 0xF9, 0xF0), line_color=ACCENT_GOLD)
add_text(slide, Inches(7.3), Inches(4.95), Inches(5.0), Inches(0.4),
         "\u6838\u5fc3\u4ef7\u503c", font_size=16, color=ACCENT_GOLD, bold=True)
insight_items = [
    "\u4f20\u7edf\u5de5\u5177\uff1a\u6bcf\u6b21\u8bca\u65ad\u4ece\u96f6\u5f00\u59cb",
    "\u672c\u5e73\u53f0\uff1a\u57fa\u4e8e\u5386\u53f2\u8ba4\u77e5\u6301\u7eed\u79ef\u7d2f",
    "\u6548\u679c\uff1a\u8bca\u65ad\u901f\u5ea6\u548c\u51c6\u786e\u7387\u968f\u65f6\u95f4\u63d0\u5347",
    "\u7c7b\u6bd4\uff1a\u4ece\"\u65b0\u6765\u7684\u5b9e\u4e60\u751f\"\u5230\"\u8001 DBA\"",
]
add_bullet_text(slide, Inches(7.3), Inches(5.35), Inches(5.0), Inches(1.6),
                insight_items, font_size=13, color=DARK_TEXT, bullet_char="\u2192")


# ============================================================
# SLIDE 15: Multi-Model + Rule Fallback
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "\u591a\u6a21\u578b\u9002\u914d + \u89c4\u5219\u5f15\u64ce\u515c\u5e95",
         font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
         "\u6a21\u578b\u65e0\u5173\u67b6\u6784", font_size=22, color=ACCENT_BLUE, bold=True)
model_items = [
    "\u7edf\u4e00 Provider \u63a5\u53e3\uff0c\u6362\u6a21\u578b\u53ea\u6539\u914d\u7f6e\u4e0d\u6539\u4ee3\u7801",
    "\u652f\u6301 OpenAI \u517c\u5bb9\u534f\u8bae\u7684\u6240\u6709\u6a21\u578b",
    "\u652f\u6301 Ollama \u672c\u5730\u90e8\u7f72\u7684\u5f00\u6e90\u6a21\u578b",
    "\u5df2\u9a8c\u8bc1\uff1a\u901a\u4e49\u5343\u95ee / DeepSeek / Kimi / ChatGPT \u7b49",
    "\u652f\u6301\u4f1a\u8bdd\u4e2d\u70ed\u5207\u6362\u6a21\u578b\uff08/model \u6307\u4ee4\uff09",
    "\u81ea\u52a8\u9002\u914d\u6a21\u578b\u80fd\u529b\u7b49\u7ea7\uff08\u5c0f\u6a21\u578b/\u5927\u6a21\u578b\uff09",
    "\u6d41\u5f0f\u8f93\u51fa + \u975e\u6d41\u5f0f\u81ea\u52a8\u964d\u7ea7",
]
add_bullet_text(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.5),
                model_items, font_size=14, color=DARK_TEXT)

add_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
         "265+ \u89c4\u5219\u5f15\u64ce\u515c\u5e95", font_size=22, color=ACCENT_GOLD, bold=True)
rule_items = [
    "\u5927\u6a21\u578b\u4e0d\u53ef\u7528\u65f6\u81ea\u52a8\u964d\u7ea7\u4e3a\u89c4\u5219\u8bca\u65ad",
    "\u89c4\u5219\u5c42\u53ea\u9648\u8ff0\u4e8b\u5b9e\uff0c\u4e0d\u8fc7\u5ea6\u63a8\u65ad",
    "\u8986\u76d6 14 \u7c7b\u6839\u56e0\u573a\u666f",
    "\u51b3\u7b56\u6811\u591a\u6b65\u5206\u6790 + SQL \u9a8c\u8bc1",
    "Oracle 265 \u6761 / PG 100+ \u6761 / MySQL \u89c4\u5219",
    "JSON \u683c\u5f0f\u89c4\u5219\u5b9a\u4e49\uff0c\u53ef\u5ba1\u8ba1\u53ef\u6269\u5c55",
    "\u65e0\u7f51\u7edc\u4f9d\u8d56\uff0c\u7eaf\u672c\u5730\u6267\u884c",
]
add_bullet_text(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.5),
                rule_items, font_size=14, color=DARK_TEXT)

add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.3),
                 CARD_BG, line_color=RGBColor(0xE0, 0xE0, 0xE0))
add_text(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.4),
         "\u667a\u80fd\u964d\u7ea7\u7b56\u7565", font_size=16, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
flow_text = ("Sentinel \u5f02\u5e38\u68c0\u6d4b  \u2192  \u4f18\u5148\u8c03\u7528\u5927\u6a21\u578b\u591a\u8f6e\u63a8\u7406  \u2192  "
             "\u82e5\u6a21\u578b\u4e0d\u53ef\u7528/\u8d85\u65f6  \u2192  \u81ea\u52a8\u964d\u7ea7\u4e3a\u89c4\u5219\u5f15\u64ce\u8bca\u65ad  \u2192  "
             "\u8f93\u51fa\u4e8b\u5b9e\u9648\u8ff0 + \u5efa\u8bae\u65b9\u6848  \u2192  \u65e0\u7f1d\u4f53\u9a8c\uff0c\u96f6\u4e2d\u65ad")
add_text(slide, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.6),
         flow_text, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 16: Deployment & Security
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "\u91d1\u878d\u7ea7\u90e8\u7f72\u65b9\u6848\uff1a\u672c\u5730\u5316 \u00b7 \u5b89\u5168 \u00b7 \u96f6\u4f9d\u8d56",
         font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

dep_cards = [
    {
        "title": "\u5355\u4e8c\u8fdb\u5236\u90e8\u7f72",
        "color": ACCENT_BLUE,
        "items": ["\u5355\u4e2a\u53ef\u6267\u884c\u6587\u4ef6\uff0c\u62f7\u8d1d\u5373\u7528",
                  "\u65e0\u9700\u5b89\u88c5\u8fd0\u884c\u65f6\u73af\u5883",
                  "\u65e0\u9700 Oracle Client \u7b49\u4f9d\u8d56",
                  "\u652f\u6301 Linux / macOS / Windows",
                  "10 \u79d2\u5b8c\u6210\u90e8\u7f72"]
    },
    {
        "title": "\u6570\u636e\u4e0d\u51fa\u57df",
        "color": RGBColor(0xCC, 0x33, 0x33),
        "items": ["\u65e0\u9700\u8fde\u63a5\u516c\u7f51\u5373\u53ef\u8fd0\u884c",
                  "\u5927\u6a21\u578b\u53ef\u672c\u5730\u5316\u90e8\u7f72\uff08Ollama\uff09",
                  "\u89c4\u5219\u5f15\u64ce\u7eaf\u672c\u5730\u6267\u884c",
                  "\u6240\u6709\u6570\u636e\u7559\u5728\u751f\u4ea7\u7f51\u7edc\u5185",
                  "\u6ee1\u8db3\u91d1\u878d\u76d1\u7ba1\u5408\u89c4\u8981\u6c42"]
    },
    {
        "title": "\u5b89\u5168\u7ba1\u63a7",
        "color": RGBColor(0x2E, 0xCC, 0x71),
        "items": ["\u56db\u7ea7\u64cd\u4f5c\u5b89\u5168\u5206\u7ea7",
                  "\u5371\u9669\u64cd\u4f5c\u5f3a\u5236\u786e\u8ba4\u4e0d\u53ef\u7ed5\u8fc7",
                  "\u5bc6\u7801\u52a0\u5bc6\u5b58\u50a8",
                  "\u53c2\u6570\u5316\u67e5\u8be2\u9632\u6ce8\u5165",
                  "\u5168\u64cd\u4f5c\u5ba1\u8ba1\u65e5\u5fd7"]
    },
    {
        "title": "\u4f01\u4e1a\u7ea7\u7279\u6027",
        "color": ACCENT_GOLD,
        "items": ["\u6279\u91cf\u6a21\u5f0f\u652f\u6301\u81ea\u52a8\u5316\u96c6\u6210",
                  "\u5b9a\u65f6\u4efb\u52a1\u5f15\u64ce\uff08Scheduler\uff09",
                  "\u591a\u5b9e\u4f8b\u7ba1\u7406\u4e0e\u5207\u6362",
                  "License \u6388\u6743\u7ba1\u7406",
                  "\u7ec8\u7aef\u72b6\u6001\u4fdd\u62a4\u4e0e\u6062\u590d"]
    },
]

for i, card in enumerate(dep_cards):
    left = Inches(0.6 + i * 3.15)
    top = Inches(1.5)
    add_rounded_rect(slide, left, top, Inches(2.9), Inches(4.5),
                     CARD_BG, line_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, left, top, Inches(2.9), Inches(0.06), card["color"])
    add_text(slide, left + Inches(0.3), top + Inches(0.25), Inches(2.3), Inches(0.4),
             card["title"], font_size=18, color=card["color"], bold=True)
    add_bullet_text(slide, left + Inches(0.3), top + Inches(0.8), Inches(2.3), Inches(3.5),
                    card["items"], font_size=13, color=DARK_TEXT, bullet_char="\u25b8")

add_rounded_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(1.0),
                 RGBColor(0xFE, 0xF9, 0xF0), line_color=ACCENT_GOLD)
deploy_flow = "\u5178\u578b\u90e8\u7f72\uff1a  \u62f7\u8d1d\u4e8c\u8fdb\u5236  \u2192  \u914d\u7f6e\u6570\u636e\u5e93\u8fde\u63a5  \u2192  \uff08\u53ef\u9009\uff09\u914d\u7f6e\u672c\u5730\u5927\u6a21\u578b  \u2192  \u5373\u523b\u4f7f\u7528"
add_text(slide, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.5),
         deploy_flow, font_size=16, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 17: Multi-Database Support
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "\u56db\u5e93\u7edf\u4e00\uff1a\u4e00\u5957\u5e73\u53f0\u8986\u76d6\u4e3b\u6d41\u6570\u636e\u5e93",
         font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

dbs = [
    {
        "name": "Oracle",
        "color": RGBColor(0xCC, 0x33, 0x33),
        "maturity": "\u6210\u719f\u5ea6: \u2605\u2605\u2605\u2605\u2605",
        "features": ["265+ \u8bca\u65ad\u89c4\u5219", "56+ \u76d1\u63a7/\u5206\u6790\u5de5\u5177", "AWR / ASH / SQL Advisor",
                     "RAC / DG \u9ad8\u53ef\u7528\u76d1\u63a7", "\u7eaf Go \u9a71\u52a8\uff0c\u65e0\u9700 Oracle Client"]
    },
    {
        "name": "PostgreSQL",
        "color": RGBColor(0x33, 0x66, 0x91),
        "maturity": "\u6210\u719f\u5ea6: \u2605\u2605\u2605\u2605",
        "features": ["100+ \u8bca\u65ad\u89c4\u5219", "\u7b49\u6548\u76d1\u63a7\u5de5\u5177\u8986\u76d6", "pg_stat \u4f53\u7cfb\u6df1\u5ea6\u96c6\u6210",
                     "\u7d22\u5f15\u4f18\u5316\u5efa\u8bae", "\u6027\u80fd\u6d1e\u5bdf\u5206\u6790"]
    },
    {
        "name": "MySQL",
        "color": RGBColor(0x00, 0x75, 0x8F),
        "maturity": "\u6210\u719f\u5ea6: \u2605\u2605\u2605\u2605",
        "features": ["Performance Schema \u96c6\u6210", "InnoDB \u5b58\u50a8\u5f15\u64ce\u5206\u6790", "\u6162\u67e5\u8be2\u65e5\u5fd7\u5206\u6790",
                     "\u9501\u4e89\u7528\u68c0\u6d4b", "\u4e3b\u4ece\u590d\u5236\u76d1\u63a7"]
    },
    {
        "name": "OpenGauss",
        "color": RGBColor(0x6C, 0x5C, 0xE7),
        "maturity": "\u6210\u719f\u5ea6: \u2605\u2605\u2605",
        "features": ["PG \u517c\u5bb9\u5185\u6838\u652f\u6301", "\u96c6\u7fa4\u7279\u6027\u76d1\u63a7", "\u7b49\u6548\u8bca\u65ad\u89c4\u5219\u8986\u76d6",
                     "\u4fe1\u521b\u9002\u914d\u9a8c\u8bc1", "\u6301\u7eed\u529f\u80fd\u589e\u5f3a\u4e2d"]
    },
]

for i, db in enumerate(dbs):
    left = Inches(0.5 + i * 3.2)
    top = Inches(1.5)
    add_rounded_rect(slide, left, top, Inches(2.95), Inches(5.0),
                     CARD_BG, line_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, left, top, Inches(2.95), Inches(0.7), db["color"])
    add_text(slide, left, top + Inches(0.1), Inches(2.95), Inches(0.5),
             db["name"], font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.2), top + Inches(0.85), Inches(2.55), Inches(0.4),
             db["maturity"], font_size=14, color=db["color"], bold=True)
    add_bullet_text(slide, left + Inches(0.2), top + Inches(1.3), Inches(2.55), Inches(3.5),
                    db["features"], font_size=13, color=DARK_TEXT, bullet_char="\u25b8")

add_text(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
         "\u7edf\u4e00\u4ea4\u4e92\u4f53\u9a8c\uff1a\u540c\u4e00\u5957 /\u6307\u4ee4 + \u81ea\u7136\u8bed\u8a00\uff0c\u8de8\u5e93\u65e0\u7f1d\u5207\u6362",
         font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 18: Value Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
         "\u6838\u5fc3\u4ef7\u503c\u603b\u7ed3", font_size=32, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(1.15), Inches(1.7), Pt(3), ACCENT_GOLD)

values = [
    ("\u964d\u672c", "\u51cf\u5c11\u5bf9\u4e13\u5bb6\u7ecf\u9a8c\u7684\u4f9d\u8d56", "\u5927\u6a21\u578b + \u89c4\u5219\u5f15\u64ce\u6301\u7eed\u79ef\u7d2f\u7ec4\u7ec7\u77e5\u8bc6\uff0c\u65b0\u624b\u4e5f\u80fd\u5feb\u901f\u5b9a\u4f4d\u95ee\u9898", ACCENT_BLUE),
    ("\u589e\u6548", "\u4ece\u544a\u8b66\u5230\u6839\u56e0\u5206\u949f\u7ea7\u95ed\u73af", "Sentinel 7\u00d724 \u76d1\u6d4b + \u81ea\u52a8\u89e6\u53d1\u591a\u8f6e\u63a8\u7406\uff0c\u65e0\u9700\u4eba\u5de5\u4ecb\u5165", RGBColor(0x2E, 0xCC, 0x71)),
    ("\u5b89\u5168", "\u6ee1\u8db3\u91d1\u878d\u76d1\u7ba1\u5408\u89c4", "\u672c\u5730\u5316\u90e8\u7f72\u3001\u6570\u636e\u4e0d\u51fa\u57df\u3001\u56db\u7ea7\u5b89\u5168\u7ba1\u63a7\u3001\u5168\u64cd\u4f5c\u5ba1\u8ba1", RGBColor(0xCC, 0x33, 0x33)),
    ("\u7edf\u4e00", "\u4e00\u5957\u5e73\u53f0\u8986\u76d6\u5f02\u6784\u5e93", "Oracle/MySQL/PG/OG \u7edf\u4e00\u4ea4\u4e92\uff0c\u6d88\u9664\u5de5\u5177\u788e\u7247\u5316", RGBColor(0x9B, 0x59, 0xB6)),
    ("\u667a\u80fd", "\u8d8a\u7528\u8d8a\u61c2\u4f60\u7684\u6570\u636e\u5e93", "\u4e0a\u4e0b\u6587\u8bb0\u5fc6\u6301\u7eed\u5b66\u4e60\uff0c\u8bca\u65ad\u51c6\u786e\u7387\u548c\u901f\u5ea6\u968f\u65f6\u95f4\u63d0\u5347", ACCENT_GOLD),
    ("\u53ef\u9760", "\u6709\u65e0\u5927\u6a21\u578b\u90fd\u80fd\u5de5\u4f5c", "\u6a21\u578b\u4e0d\u53ef\u7528\u81ea\u52a8\u964d\u7ea7\u89c4\u5219\u5f15\u64ce\uff0c\u96f6\u4e2d\u65ad\u4fdd\u969c", RGBColor(0x34, 0x49, 0x5E)),
]

for i, (title, subtitle, desc, color) in enumerate(values):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.15)
    top = Inches(1.5 + row * 2.8)
    add_rounded_rect(slide, left, top, Inches(3.85), Inches(2.5),
                     CARD_BG, line_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, left, top, Inches(3.85), Inches(0.06), color)
    add_text(slide, left + Inches(0.3), top + Inches(0.2), Inches(3.25), Inches(0.4),
             title, font_size=24, color=color, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.65), Inches(3.25), Inches(0.35),
             subtitle, font_size=14, color=DARK_TEXT, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(1.1), Inches(3.25), Inches(1.2),
             desc, font_size=13, color=MID_GRAY)


# ============================================================
# SLIDE 19: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT_GOLD)
add_text(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.2),
         "\u8c22\u8c22", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.5), Inches(3.8), Inches(2.3), Pt(2), ACCENT_GOLD)
add_text(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
         PRODUCT_NAME, font_size=24, color=SUBTITLE_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
         "\u9762\u5411\u91d1\u878d\u7ea7\u751f\u4ea7\u73af\u5883\u7684\u6570\u636e\u5e93\u667a\u80fd\u8fd0\u7ef4\u89e3\u51b3\u65b9\u6848",
         font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = "/Users/yingjiewang/opendb/\u6570\u636e\u5e93\u8fd0\u7ef4\u667a\u80fd\u5316\u5e73\u53f0\u4ecb\u7ecd.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
